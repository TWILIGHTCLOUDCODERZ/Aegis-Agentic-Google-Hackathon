# Build Aegis pitch deck as a native, editable .pptx (16:9), dark theme.
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- palette ----
BG      = RGBColor(0x07, 0x0F, 0x1C)
BG2     = RGBColor(0x0B, 0x18, 0x2C)
PANEL   = RGBColor(0x10, 0x22, 0x39)
LINE    = RGBColor(0x2A, 0x40, 0x59)
INK     = RGBColor(0xEE, 0xF4, 0xFB)
MUTED   = RGBColor(0xA2, 0xB8, 0xCE)
FAINT   = RGBColor(0x6D, 0x87, 0xA0)
BLUE    = RGBColor(0x42, 0x85, 0xF4)
BLUELT  = RGBColor(0x78, 0xA8, 0xFF)
GREEN   = RGBColor(0x5C, 0xD0, 0x8B)
YELLOW  = RGBColor(0xFB, 0xBC, 0x04)
RED     = RGBColor(0xFF, 0x8A, 0x80)
PURPLE  = RGBColor(0xC4, 0xB0, 0xFF)

DISPLAY = "Segoe UI Semibold"   # widely available; stands in for Space Grotesk
BODY    = "Segoe UI"
MONO    = "Consolas"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height

LOGO = r"D:\DEVPOST\google hackathon\project\pitch\team-tcc-logo.png"

def slide():
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = BG; r.line.fill.background()
    r.shadow.inherit = False
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2, r._element)
    # Team TCC logo — top-left corner of every slide
    s.shapes.add_picture(LOGO, Inches(0.34), Inches(0.32), height=Inches(0.66))
    return s

def _noshadow(sp): sp.shadow.inherit = False

def box(s, x, y, w, h, fill=None, line=None, line_w=1.0, radius=0.08):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    try: shp.adjustments[0] = radius
    except Exception: pass
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    _noshadow(shp)
    return shp

def dot(s, x, y, d, color):
    shp = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    shp.fill.solid(); shp.fill.fore_color.rgb = color; shp.line.fill.background(); _noshadow(shp)
    return shp

def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=6, line=1.12):
    """runs: list of paragraphs; each paragraph = list of (txt, size, color, font, bold, spacing)"""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space_after); p.space_before = Pt(0)
        try: p.line_spacing = line
        except Exception: pass
        for (txt, size, color, font, bold, *rest) in para:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.color.rgb = color
            r.font.name = font; r.font.bold = bold
            if rest:
                spc = rest[0]
                rPr = r._r.get_or_add_rPr(); rPr.set('spc', str(int(spc*100)))
    return tb

def eyebrow(s, num, label, sub=""):
    dot(s, 0.92, 0.99, 0.11, BLUE)
    parts = [(num+"  ", 12, FAINT, MONO, False, 2), (label, 12, BLUELT, MONO, True, 2)]
    if sub: parts.append(("   / "+sub, 12, FAINT, MONO, False, 2))
    text(s, 1.15, 0.86, 11, 0.5, [parts])

def title(s, t, size=34, y=1.35, color=INK, w=11.5):
    text(s, 0.9, y, w, 1.4, [[(t, size, color, DISPLAY, True)]], line=1.02)

def card(s, x, y, w, h, tag, head, body, head_color=INK, fill=PANEL, line=LINE):
    box(s, x, y, w, h, fill=fill, line=line, line_w=1.0)
    pad = 0.22
    text(s, x+pad, y+pad, w-2*pad, h-2*pad, [
        [(tag, 10, FAINT, MONO, True, 1.5)],
        [(head, 15, head_color, DISPLAY, True)],
        [(body, 11.5, MUTED, BODY, False)],
    ], space_after=6, line=1.16)

def metric(s, x, y, w, n, nsuf, label, color, sub=""):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(1.9))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left=tf.margin_right=tf.margin_top=tf.margin_bottom=0
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = n; r.font.size=Pt(46); r.font.bold=True; r.font.name=DISPLAY; r.font.color.rgb=color
    if nsuf:
        r2 = p.add_run(); r2.text=" "+nsuf; r2.font.size=Pt(20); r2.font.bold=True; r2.font.name=DISPLAY; r2.font.color.rgb=FAINT
    p2 = tf.add_paragraph(); p2.space_before=Pt(6); r=p2.add_run(); r.text=label; r.font.size=Pt(13); r.font.name=BODY; r.font.color.rgb=MUTED
    if sub:
        p3 = tf.add_paragraph(); p3.space_before=Pt(3); r=p3.add_run(); r.text=sub; r.font.size=Pt(11); r.font.name=MONO; r.font.color.rgb=FAINT

def chip(s, x, y, w, label, color):
    h=0.42
    box(s, x, y, w, h, fill=PANEL, line=LINE, line_w=1.0, radius=0.5)
    dot(s, x+0.16, y+h/2-0.055, 0.11, color)
    text(s, x+0.38, y, w-0.4, h, [[(label, 11, MUTED, MONO, False)]], anchor=MSO_ANCHOR.MIDDLE)

# ============================================================ 1 · TITLE
s = slide()
# brand mark
bm = box(s, 0.9, 0.78, 0.72, 0.72, fill=BG2, line=LINE, line_w=1.2, radius=0.24)
dot(s, 1.06, 0.95, 0.38, RGBColor(0x14,0x30,0x5a))
text(s, 0.95, 0.86, 0.62, 0.55, [[("A", 26, BLUELT, DISPLAY, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, 1.78, 0.80, 8, 0.8, [
    [("AEGIS", 22, INK, DISPLAY, True, 3)],
    [("AI FRAUD INTELLIGENCE", 10, FAINT, MONO, False, 3)],
], space_after=2)
text(s, 0.88, 2.35, 11.5, 2.6, [
    [("Fraud decisions in", 62, INK, DISPLAY, True)],
    [("milliseconds.", 62, BLUELT, DISPLAY, True)],
], line=0.98, space_after=0)
text(s, 0.9, 5.05, 9.6, 1.0, [[(
    "An autonomous, real-time fraud & financial-crime defense platform for banking — built on Gemini 3.5 and Google Cloud.",
    18, MUTED, BODY, False)]], line=1.35)
for i,(lbl,c) in enumerate([("Multi-agent investigation",BLUE),("Google Cloud · Gemini 3.5",GREEN),("All Things Agentic · Taskmaster",YELLOW)]):
    chip(s, 0.9+i*3.35, 6.4, 3.15, lbl, c)

# ============================================================ 2 · OVERVIEW
s = slide()
eyebrow(s, "01", "OVERVIEW", "the banking sector")
title(s, "Every transaction is a 100-millisecond decision.", 33)
text(s, 0.9, 2.5, 11.3, 1.4, [[(
    "Modern banks and card networks approve or decline thousands of payments every second. Behind each swipe sits one unforgiving question — is this really the customer? — that must be answered instantly, at scale, with money and trust on the line.",
    18, MUTED, BODY, False)]], line=1.4)
data=[("100 ms","","Budget to decide on a live authorization"),
      ("$1T+","","Digital payments flowing through banks yearly"),
      ("24 / 7","","Attack surface — fraud never sleeps")]
for i,(n,suf,l) in enumerate(data):
    metric(s, 0.95+i*3.9, 4.55, 3.6, n, suf, l, BLUELT)

# ============================================================ 3 · PROBLEM
s = slide()
eyebrow(s, "02", "THE PROBLEM")
title(s, "Static rules force an impossible trade-off.", 33)
text(s, 0.9, 2.42, 11.3, 1.2, [[(
    "Legacy fraud engines are frozen if-then rules. Tighten them and you decline good customers; loosen them and fraud slips through. Either way a human reviews the aftermath hours later — long after the money is gone.",
    16, MUTED, BODY, False)]], line=1.35)
cards=[("Fraud slips through","Rules only catch patterns someone already wrote down. Novel attacks pass untouched.",RED),
       ("False declines","Good customers blocked at checkout — the biggest driver of churn and lost revenue.",YELLOW),
       ("Slow to act","Manual review comes hours later. Explanations reconstructed after the fact.",BLUELT),
       ("No memory","Every transaction judged cold. The system never learns this customer is legit.",MUTED)]
cw=2.86; gap=0.18; x0=0.9
for i,(h,b,c) in enumerate(cards):
    card(s, x0+i*(cw+gap), 4.35, cw, 2.5, "", h, b, head_color=c)

# ============================================================ 4 · SOLUTION
s = slide()
eyebrow(s, "03", "THE AEGIS SOLUTION")
title(s, "A team of AI agents that investigates every risky payment.", 30, w=12)
text(s, 0.9, 2.5, 11.4, 1.3, [[(
    "Aegis replaces the static rulebook with an autonomous investigation. Low-risk payments clear in milliseconds; anything suspicious goes to specialist Gemini 3.5 agents that gather evidence, weigh it, and return a decision with its reasons — then verify the customer instead of declining them.",
    16, MUTED, BODY, False)]], line=1.35)
cards=[("AUTONOMOUS","Investigates, not just scores","Agents pull evidence, debate, and reach a verdict — no analyst in the loop."),
       ("ADAPTIVE","Verify, don't decline","Borderline? Step up with an OTP instead of losing a good customer."),
       ("EXPLAINABLE","Reason codes + evidence","Every decision ships with confidence, reason codes and an evidence summary."),
       ("REMEMBERS","Customer memory","Firestore memory of confirmed-legit behavior means fewer repeat challenges.")]
for i,(tg,h,b) in enumerate(cards):
    card(s, x0+i*(cw+gap), 4.5, cw, 2.4, tg, h, b)

# ============================================================ 5 · HOW IT WORKS
s = slide()
eyebrow(s, "04", "HOW IT WORKS", "two-speed architecture")
title(s, "Fast enough for checkout. Deep enough to catch fraud.", 30, w=12)
# two lanes
lane_y=2.5; lane_h=3.2; lane_w=5.7
box(s, 0.9, lane_y, lane_w, lane_h, fill=RGBColor(0x0E,0x1F,0x38), line=RGBColor(0x27,0x45,0x6b), line_w=1.2, radius=0.06)
box(s, 6.9, lane_y, lane_w, lane_h, fill=RGBColor(0x16,0x14,0x2e), line=RGBColor(0x3a,0x30,0x60), line_w=1.2, radius=0.06)
text(s, 1.15, lane_y+0.18, 4, 0.4, [[("FAST PATH", 15, INK, DISPLAY, True)]])
text(s, 1.15, lane_y+0.18, lane_w-0.5, 0.4, [[("◈ milliseconds", 11, BLUELT, MONO, False)]], align=PP_ALIGN.RIGHT)
text(s, 7.15, lane_y+0.18, 4, 0.4, [[("DEEP PATH", 15, INK, DISPLAY, True)]])
text(s, 7.15, lane_y+0.18, lane_w-0.5, 0.4, [[("◈ async · multi-agent", 11, PURPLE, MONO, False)]], align=PP_ALIGN.RIGHT)
fast=[("01","Transaction arrives","Validated & enriched, streamed over Pub/Sub."),
      ("02","Rules + Gemini Flash pre-filter","Instant triage on the live authorization."),
      ("03","Clear or escalate","Low risk → approve now. High risk → deep path.")]
deep=[("04","Orchestrator convenes the team","Six specialist agents investigate on Gemini 3.5."),
      ("05","Evidence + critic review","Findings cross-checked; a critic challenges the verdict."),
      ("06","Adaptive decision","Approve · Step-up · Hold · Block — with reasons.")]
def steps(s, x, data):
    yy=lane_y+0.72
    for b,t,d in data:
        text(s, x, yy, 0.4, 0.3, [[(b, 11, FAINT, MONO, True)]])
        text(s, x+0.42, yy, 4.9, 0.8, [[(t, 13, INK, DISPLAY, True)],[(d, 11.5, MUTED, BODY, False)]], space_after=2, line=1.1)
        yy+=0.78
steps(s, 1.2, fast); steps(s, 7.2, deep)
# decision chips
labels=[("Approve",GREEN),("Step-up · OTP",YELLOW),("Hold · review",BLUELT),("Block",RED)]
xx=0.9
for lbl,c in labels:
    w=1.55 if len(lbl)<8 else 2.0
    chip(s, xx, 6.05, w, lbl, c); xx+=w+0.15
text(s, xx+0.1, 6.05, 4.8, 0.42, [[("every verdict carries reason codes · confidence · evidence", 10.5, FAINT, MONO, False)]], anchor=MSO_ANCHOR.MIDDLE)

# ============================================================ 6 · TYSON
s = slide()
eyebrow(s, "05", "TYSON USER FLOW", "step-up in action")
title(s, "The same payment. Two very different outcomes.", 30, w=12)
# scene bar
box(s, 0.9, 2.45, 11.5, 0.95, fill=PANEL, line=LINE, line_w=1.0, radius=0.12)
dot(s, 1.12, 2.68, 0.5, BG2)
text(s, 1.12, 2.68, 0.5, 0.5, [[("T", 20, BLUELT, DISPLAY, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, 1.8, 2.55, 10.4, 0.8, [
    [("Tyson · large auto-pay charge", 15, INK, DISPLAY, True)],
    [("A high-value payment fires. Aegis checks card status, then the step-up policy gate — before it ever reaches the agents.", 12.5, MUTED, BODY, False)],
], space_after=2, line=1.1)
# two paths
pw=5.7; py=3.7; ph=2.15
box(s, 0.9, py, pw, ph, fill=RGBColor(0x22,0x11,0x12), line=RGBColor(0x6a,0x2a,0x28), line_w=1.2, radius=0.07)
box(s, 6.9, py, pw, ph, fill=RGBColor(0x0f,0x22,0x18), line=RGBColor(0x27,0x55,0x38), line_w=1.2, radius=0.07)
text(s, 1.15, py+0.2, pw-0.5, 0.35, [[("◆ auto-pay OFF · no OTP verified", 12, RED, MONO, False)]])
text(s, 1.15, py+0.55, pw-0.5, 0.5, [[("✕  Card blocked", 24, RED, DISPLAY, True)]])
text(s, 1.15, py+1.15, pw-0.5, 0.9, [[("Unverified high-value charge with auto-pay disabled — Aegis refuses it and protects the customer, with a reason code the analyst sees instantly.", 12, MUTED, BODY, False)]], line=1.2)
text(s, 7.15, py+0.2, pw-0.5, 0.35, [[("◆ OTP verified / auto-pay ON", 12, GREEN, MONO, False)]])
text(s, 7.15, py+0.55, pw-0.5, 0.5, [[("✓  Approved", 24, GREEN, DISPLAY, True)]])
text(s, 7.15, py+1.15, pw-0.5, 0.9, [[("Tyson confirms with a Mobile or Email OTP — Aegis verifies identity and approves the same payment. A good customer is kept, not declined.", 12, MUTED, BODY, False)]], line=1.2)
text(s, 0.9, 6.15, 11.5, 0.6, [[("The lesson: ", 15, MUTED, BODY, False),("verify the customer, don't just decline the charge.", 15, INK, DISPLAY, True)]])

# ============================================================ 7 · AGENTS
s = slide()
eyebrow(s, "05", "BEHIND THE SCENES", "how the agents work")
title(s, "An investigation team, not a single model call.", 30, w=12)
text(s, 0.9, 2.42, 11.4, 0.9, [[(
    "When the deep path fires, an Orchestrator convenes specialists on Gemini 3.5. Each has one job; a Critic keeps them honest. Every step streams live in the console.",
    15, MUTED, BODY, False)]], line=1.3)
agents=[("Orchestrator","Convenes the team, routes evidence, issues the final decision.",BLUE,True),
        ("Card Status","Checks stolen / lost / frozen — a hard block short-circuits all.",PURPLE,False),
        ("Step-Up Control","Applies the OTP & auto-pay policy gate — verify vs decline.",PURPLE,False),
        ("Investigator","Reconstructs the story: amount, location, device, history.",PURPLE,False),
        ("Network Analyst","Looks for mule rings and connected-account patterns.",PURPLE,False),
        ("Intel","Cross-references fraud signals, watchlists and threat intel.",PURPLE,False),
        ("Compliance","Frames the case for AML / SAR; drafts the narrative.",PURPLE,False),
        ("Critic","Challenges the verdict and stress-tests confidence.",YELLOW,False)]
aw=2.86; ah=1.55; ax=0.9; ay=3.45
for i,(n,d,c,orch) in enumerate(agents):
    col=i%4; row=i//4
    x=ax+col*(aw+0.18); y=ay+row*(ah+0.18)
    fill=RGBColor(0x11,0x22,0x3c) if orch else PANEL
    ln=RGBColor(0x27,0x45,0x6b) if orch else LINE
    box(s, x, y, aw, ah, fill=fill, line=ln, line_w=1.0, radius=0.09)
    sq=s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x+0.2), Inches(y+0.24), Inches(0.13), Inches(0.13))
    sq.fill.solid(); sq.fill.fore_color.rgb=c; sq.line.fill.background(); _noshadow(sq)
    text(s, x+0.42, y+0.16, aw-0.6, 0.35, [[(n, 14, INK, DISPLAY, True)]])
    text(s, x+0.2, y+0.6, aw-0.4, 0.9, [[(d, 11, MUTED, BODY, False)]], line=1.15)

# ============================================================ 8 · RESULTS
s = slide()
eyebrow(s, "05", "BEHIND THE SCENES", "measured impact")
title(s, "Benchmarked against a rules-only engine.", 32)
text(s, 0.9, 2.45, 11.3, 0.8, [[(
    "Aegis and a traditional rules baseline were run over the same labeled transaction set. The gap is the whole point.",
    16, MUTED, BODY, False)]], line=1.3)
res=[("95.2","%","Fraud caught","rules baseline: 12.7%",GREEN),
     ("90.4","%","Fewer false declines","39.0% → 3.8%",BLUELT),
     ("88.2","%","Block precision","when Aegis blocks, it's right",PURPLE),
     ("<100","ms","Fast-path decision","deep path runs async",YELLOW)]
for i,(n,suf,l,sub,c) in enumerate(res):
    metric(s, 0.95+i*3.05, 3.7, 2.9, n, suf, l, c, sub)
text(s, 0.9, 6.35, 11.4, 0.5, [[(
    "Evaluated on a labeled synthetic transaction set; Aegis decision engine vs a static rules baseline.",
    12, FAINT, MONO, False)]])

# ============================================================ 9 · TECHNOLOGY
s = slide()
eyebrow(s, "06", "TECHNOLOGY")
title(s, "Built entirely on Google Cloud & Gemini 3.5.", 32)
def tstack(s, x, heading, rows):
    text(s, x, 2.55, 5.4, 0.35, [[(heading, 11, FAINT, MONO, True, 1.5)]])
    yy=3.05
    for k,v,sub in rows:
        box(s, x, yy, 5.5, 0.62, fill=PANEL, line=LINE, line_w=1.0, radius=0.14)
        text(s, x+0.2, yy, 1.1, 0.62, [[(k, 10.5, FAINT, MONO, True, 0.8)]], anchor=MSO_ANCHOR.MIDDLE)
        text(s, x+1.35, yy, 4.0, 0.62, [[(v+"  ", 13, INK, BODY, True),(sub, 11.5, MUTED, BODY, False)]], anchor=MSO_ANCHOR.MIDDLE)
        yy+=0.72
left=[("MODEL","Gemini 3.5","· Flash fast-path + agents"),
      ("AGENTS","Google GenAI SDK","· orchestration"),
      ("API","FastAPI on Cloud Run","· SSE streaming"),
      ("EVENTS","Pub/Sub","· transaction stream"),
      ("MEMORY","Firestore","· cases & customer memory")]
right=[("AUTH","Firebase Auth","· + guest mode"),
       ("SECRETS","Secret Manager","· keys & config"),
       ("SAFETY","Model Armor","· prompt & output guardrails"),
       ("BUILD","Cloud Build","· deploy from source"),
       ("CONSOLE","React · TypeScript","· Vite on Cloud Run")]
tstack(s, 0.9, "INTELLIGENCE & BACKEND", left)
tstack(s, 6.95, "PLATFORM & EXPERIENCE", right)

# ============================================================ 10 · CONCLUSION
s = slide()
eyebrow(s, "07", "CONCLUSION")
text(s, 0.9, 1.5, 11.2, 1.8, [
    [("Aegis turns fraud defense from a static gate", 30, INK, DISPLAY, True)],
    [("into an ", 30, INK, DISPLAY, True),("autonomous investigator", 30, BLUELT, DISPLAY, True)],
    [("that thinks, verifies, and explains — in real time.", 30, INK, DISPLAY, True)],
], line=1.12, space_after=0)
ta=[("Catch more, decline less","95% of fraud caught while cutting false declines 90% — protecting revenue and trust at once."),
    ("Verify, don't reject","Adaptive step-up keeps good customers through checkout instead of turning them away."),
    ("Cloud-native & explainable","Serverless on Google Cloud, powered by Gemini 3.5, with a reason for every decision.")]
for i,(t,d) in enumerate(ta):
    x=0.9+i*3.9
    ln=s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(4.35), Inches(3.6), Inches(0.03))
    ln.fill.solid(); ln.fill.fore_color.rgb=BLUE; ln.line.fill.background(); _noshadow(ln)
    text(s, x, 4.5, 3.6, 1.5, [[(t, 15, INK, DISPLAY, True)],[(d, 12, MUTED, BODY, False)]], space_after=4, line=1.2)
for i,(lbl,c) in enumerate([("Aegis · AI Fraud Intelligence",BLUE),("Built by Deepan",GREEN),("Thank you",YELLOW)]):
    chip(s, 0.9+i*3.6, 6.5, 3.4, lbl, c)

out = r"D:\DEVPOST\google hackathon\project\pitch\Aegis-Pitch-Deck.pptx"
prs.save(out)
print("saved", out, "slides:", len(prs.slides._sldIdLst))
